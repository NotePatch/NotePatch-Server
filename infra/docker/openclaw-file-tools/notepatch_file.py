#!/usr/bin/env python3
"""Safely inspect and normalize task-local files for NotePatch OpenClaw."""

from __future__ import annotations

import argparse
import csv
import email
import hashlib
import json
import mimetypes
import os
import shutil
import stat
import subprocess
import sys
import tarfile
import tempfile
import zipfile
from dataclasses import asdict, dataclass
from email import policy
from pathlib import Path, PurePosixPath
from typing import Any


WORKSPACE_ROOT = Path(os.getenv("OPENCLAW_WORKSPACE_ROOT", "/workspace")).resolve()
MAX_FILE_BYTES = int(os.getenv("OPENCLAW_FILE_PARSE_MAX_FILE_MB", "200")) * 1024 * 1024
MAX_OUTPUT_BYTES = int(os.getenv("OPENCLAW_FILE_PARSE_MAX_OUTPUT_MB", "20")) * 1024 * 1024
MAX_ARCHIVE_ENTRIES = int(os.getenv("OPENCLAW_ARCHIVE_MAX_ENTRIES", "1000"))
MAX_ARCHIVE_BYTES = int(os.getenv("OPENCLAW_ARCHIVE_MAX_EXPANDED_MB", "500")) * 1024 * 1024
MAX_ARCHIVE_RATIO = float(os.getenv("OPENCLAW_ARCHIVE_MAX_RATIO", "100"))
MAX_ARCHIVE_DEPTH = int(os.getenv("OPENCLAW_ARCHIVE_MAX_DEPTH", "3"))
MAX_MEDIA_AUDIO_SECONDS = int(os.getenv("OPENCLAW_MEDIA_AUDIO_MAX_SECONDS", "600"))
COMMAND_TIMEOUT = int(os.getenv("OPENCLAW_FILE_PARSE_TIMEOUT_SECONDS", "180"))


class FileToolError(RuntimeError):
    pass


@dataclass
class Inspection:
    path: str
    filename: str
    size: int
    mime_type: str
    sha256: str
    encrypted: bool
    category: str


def _within(path: Path, root: Path) -> Path:
    resolved = path.expanduser().resolve()
    if not resolved.is_relative_to(root):
        raise FileToolError(f"Path is outside {root}")
    return resolved


def source_path(value: str) -> Path:
    path = _within(Path(value), WORKSPACE_ROOT)
    if not path.is_file():
        raise FileToolError("Input file does not exist")
    size = path.stat().st_size
    if size <= 0:
        raise FileToolError("Input file is empty")
    if size > MAX_FILE_BYTES:
        raise FileToolError(f"Input exceeds {MAX_FILE_BYTES // (1024 * 1024)} MB")
    return path


def output_path(value: str) -> Path:
    path = _within(Path(value), WORKSPACE_ROOT)
    normalized = PurePosixPath(path.as_posix())
    parts = normalized.parts
    if "tasks" not in parts or "output" not in parts or "parser" not in parts:
        raise FileToolError("Output must be inside a task output/parser directory")
    path.mkdir(parents=True, exist_ok=True)
    return path


def run(command: list[str], *, cwd: Path | None = None, timeout: int = COMMAND_TIMEOUT) -> subprocess.CompletedProcess[str]:
    try:
        result = subprocess.run(
            command,
            cwd=cwd,
            text=True,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=timeout,
            check=False,
            env={**os.environ, "HOME": str(cwd or Path("/tmp"))},
        )
    except subprocess.TimeoutExpired as exc:
        raise FileToolError(f"Parser timed out after {timeout} seconds") from exc
    if result.returncode != 0:
        detail = (result.stderr or result.stdout).strip()[:500]
        raise FileToolError(f"Parser command failed: {detail or command[0]}")
    return result


def mime_type(path: Path) -> str:
    try:
        import magic

        detected = str(magic.from_file(str(path), mime=True) or "").lower()
    except Exception:
        detected = ""
    return detected or mimetypes.guess_type(path.name)[0] or "application/octet-stream"


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as stream:
        for chunk in iter(lambda: stream.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _is_zip_encrypted(path: Path) -> bool:
    if not zipfile.is_zipfile(path):
        return False
    with zipfile.ZipFile(path) as archive:
        return any(item.flag_bits & 0x1 for item in archive.infolist())


def _is_pdf_encrypted(path: Path) -> bool:
    if path.suffix.lower() != ".pdf":
        return False
    result = subprocess.run(["qpdf", "--is-encrypted", str(path)], capture_output=True, timeout=20)
    return result.returncode == 0


def _is_office_encrypted(path: Path) -> bool:
    if path.suffix.lower() not in {".doc", ".docx", ".xls", ".xlsx", ".ppt", ".pptx"}:
        return False
    try:
        import msoffcrypto

        with path.open("rb") as stream:
            return bool(msoffcrypto.OfficeFile(stream).is_encrypted())
    except Exception:
        return False


def is_encrypted(path: Path) -> bool:
    return _is_pdf_encrypted(path) or _is_zip_encrypted(path) or _is_office_encrypted(path)


def category(path: Path, detected: str) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf" or detected == "application/pdf":
        return "pdf"
    if suffix in {".doc", ".docx", ".odt", ".rtf"}:
        return "document"
    if suffix in {".ppt", ".pptx", ".odp"}:
        return "presentation"
    if suffix in {".xls", ".xlsx", ".ods"}:
        return "spreadsheet"
    if suffix in {".epub"}:
        return "ebook"
    if suffix in {".eml", ".msg"}:
        return "email"
    if suffix == ".ipynb":
        return "notebook"
    if suffix in {".zip", ".7z", ".rar", ".tar", ".tgz", ".gz", ".bz2", ".xz", ".zst"}:
        return "archive"
    if detected.startswith("image/"):
        return "image"
    if detected.startswith("audio/") or detected.startswith("video/"):
        return "media"
    if detected.startswith("text/") or suffix in {".md", ".csv", ".tsv", ".json", ".jsonl", ".yaml", ".yml", ".xml", ".html", ".htm"}:
        return "text"
    return "unknown"


def inspect_file(path: Path) -> Inspection:
    detected = mime_type(path)
    return Inspection(
        path=str(path), filename=path.name, size=path.stat().st_size, mime_type=detected,
        sha256=sha256(path), encrypted=is_encrypted(path), category=category(path, detected),
    )


def decode_text(path: Path) -> str:
    from charset_normalizer import from_bytes

    raw = path.read_bytes()
    match = from_bytes(raw).best()
    return str(match) if match is not None else raw.decode("utf-8", errors="replace")


def extract_pdf(path: Path, out: Path) -> tuple[str, str, dict[str, Any]]:
    import pymupdf

    document = pymupdf.open(path)
    pages = []
    for page_index, page in enumerate(document):
        pages.append(f"## Page {page_index + 1}\n\n{page.get_text('text').strip()}")
    text = "\n\n".join(pages)
    info_raw = run(["pdfinfo", str(path)]).stdout
    info = dict(line.split(":", 1) for line in info_raw.splitlines() if ":" in line)
    return text, text, {"pages": len(document), "pdfinfo": info}


def extract_docx(path: Path) -> tuple[str, str, dict[str, Any]]:
    from docx import Document

    document = Document(path)
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
        if rows:
            lines.extend([" | ".join(rows[0]), " | ".join("---" for _ in rows[0])])
            lines.extend(" | ".join(row) for row in rows[1:])
    text = "\n\n".join(lines)
    return text, text, {"paragraphs": len(document.paragraphs), "tables": len(document.tables)}


def extract_pptx(path: Path) -> tuple[str, str, dict[str, Any]]:
    from pptx import Presentation

    presentation = Presentation(path)
    pages: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        values = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        pages.append(f"## Slide {index}\n\n" + "\n\n".join(values))
    text = "\n\n".join(pages)
    return text, text, {"slides": len(presentation.slides)}


def extract_xlsx(path: Path) -> tuple[str, str, dict[str, Any]]:
    from openpyxl import load_workbook

    book = load_workbook(path, read_only=True, data_only=True)
    sections: list[str] = []
    for sheet in book.worksheets:
        sections.append(f"## Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                sections.append(" | ".join(values))
    text = "\n".join(sections)
    return text, text, {"sheets": book.sheetnames}


def extract_epub(path: Path) -> tuple[str, str, dict[str, Any]]:
    import ebooklib
    from bs4 import BeautifulSoup
    from ebooklib import epub

    book = epub.read_epub(str(path), options={"ignore_ncx": True})
    sections = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        content = BeautifulSoup(item.get_content(), "lxml").get_text("\n", strip=True)
        if content:
            sections.append(content)
    text = "\n\n".join(sections)
    return text, text, {"sections": len(sections)}


def extract_email(path: Path) -> tuple[str, str, dict[str, Any]]:
    if path.suffix.lower() == ".msg":
        import extract_msg

        message = extract_msg.Message(str(path))
        headers = {"subject": message.subject, "from": message.sender, "to": message.to, "date": message.date}
        body = message.body or ""
    else:
        message = email.message_from_bytes(path.read_bytes(), policy=policy.default)
        headers = {key.lower(): str(message.get(key, "")) for key in ("Subject", "From", "To", "Date")}
        body = message.get_body(preferencelist=("plain", "html"))
        body = body.get_content() if body is not None else ""
    text = "\n".join(f"{key}: {value}" for key, value in headers.items()) + "\n\n" + str(body)
    return text, text, {"headers": headers}


def extract_notebook(path: Path) -> tuple[str, str, dict[str, Any]]:
    import nbformat

    notebook = nbformat.read(path, as_version=4)
    sections = []
    for index, cell in enumerate(notebook.cells, start=1):
        sections.append(f"## Cell {index} ({cell.cell_type})\n\n{cell.source}")
        for output in cell.get("outputs", []):
            value = output.get("text") or (output.get("data") or {}).get("text/plain")
            if value:
                sections.append(str(value))
    text = "\n\n".join(sections)
    return text, text, {"cells": len(notebook.cells)}


def extract_image(path: Path, out: Path) -> tuple[str, str, dict[str, Any]]:
    from PIL import Image, ImageOps

    with Image.open(path) as image:
        metadata = {"width": image.width, "height": image.height, "format": image.format, "mode": image.mode}
        preview = ImageOps.exif_transpose(image).convert("RGB")
        preview.thumbnail((2048, 2048))
        preview.save(out / "preview.jpg", "JPEG", quality=85, optimize=True)
    languages = os.getenv("OPENCLAW_TESSERACT_LANG", "eng+chi_sim+por")
    result = subprocess.run(["tesseract", str(path), "stdout", "-l", languages], text=True, capture_output=True, timeout=COMMAND_TIMEOUT)
    text = result.stdout if result.returncode == 0 else ""
    metadata["ocr"] = "tesseract" if text.strip() else "unavailable"
    return text, text, metadata


def extract_media(path: Path, out: Path) -> tuple[str, str, dict[str, Any]]:
    probe = json.loads(run(["ffprobe", "-v", "quiet", "-print_format", "json", "-show_format", "-show_streams", str(path)]).stdout)
    streams = probe.get("streams", [])
    if any(item.get("codec_type") == "video" for item in streams):
        previews = out / "previews"
        previews.mkdir(exist_ok=True)
        run(["ffmpeg", "-v", "error", "-i", str(path), "-vf", "fps=1/30,scale=1280:-2", "-frames:v", "6", str(previews / "frame-%02d.jpg")])
    audio_output = None
    if any(item.get("codec_type") == "audio" for item in streams):
        audio_output = out / "audio-preview.m4a"
        run([
            "ffmpeg", "-v", "error", "-i", str(path), "-vn", "-t",
            str(MAX_MEDIA_AUDIO_SECONDS), "-c:a", "aac", "-b:a", "96k", str(audio_output),
        ])
    summary = json.dumps(probe, ensure_ascii=False, indent=2)
    return summary, summary, {
        "ffprobe": probe,
        "audio_preview": audio_output.name if audio_output is not None else None,
        "audio_preview_max_seconds": MAX_MEDIA_AUDIO_SECONDS,
        "transcription": "not_available",
    }


def _safe_member(name: str) -> PurePosixPath:
    value = PurePosixPath(name.replace("\\", "/"))
    if value.is_absolute() or ".." in value.parts:
        raise FileToolError(f"Unsafe archive member: {name}")
    return value


def archive_listing(path: Path) -> list[dict[str, Any]]:
    entries: list[dict[str, Any]] = []
    if zipfile.is_zipfile(path):
        with zipfile.ZipFile(path) as archive:
            for item in archive.infolist():
                _safe_member(item.filename)
                if stat.S_ISLNK(item.external_attr >> 16):
                    raise FileToolError("Archive contains a symbolic link")
                entries.append({"name": item.filename, "size": item.file_size, "compressed": item.compress_size})
    elif tarfile.is_tarfile(path):
        with tarfile.open(path) as archive:
            for item in archive.getmembers():
                _safe_member(item.name)
                if item.issym() or item.islnk() or item.isdev():
                    raise FileToolError("Archive contains a link or device")
                entries.append({"name": item.name, "size": item.size, "compressed": 0})
    else:
        listing = run(["7z", "l", "-slt", str(path)]).stdout
        current: dict[str, str] = {}
        for line in listing.splitlines() + [""]:
            if " = " in line:
                key, value = line.split(" = ", 1)
                current[key] = value
            elif current:
                name = current.get("Path")
                if name and name != str(path):
                    _safe_member(name)
                    if current.get("Encrypted") == "+":
                        raise FileToolError("Password-protected archives are not supported")
                    attributes = (current.get("Attributes") or "").lower()
                    mode = (current.get("Mode") or "").lower()
                    if attributes.startswith("l") or mode.startswith("l"):
                        raise FileToolError("Archive contains a symbolic link")
                    entries.append({"name": name, "size": int(current.get("Size") or 0), "compressed": int(current.get("Packed Size") or 0)})
                current = {}
    if len(entries) > MAX_ARCHIVE_ENTRIES:
        raise FileToolError(f"Archive exceeds {MAX_ARCHIVE_ENTRIES} entries")
    expanded = sum(int(item["size"]) for item in entries)
    compressed = max(sum(int(item["compressed"]) for item in entries), 1)
    if expanded > MAX_ARCHIVE_BYTES or expanded / compressed > MAX_ARCHIVE_RATIO:
        raise FileToolError("Archive expansion limit exceeded")
    return entries


def extract_archive(path: Path, out: Path) -> tuple[str, str, dict[str, Any]]:
    if sum(part == "extracted" for part in path.parts) >= MAX_ARCHIVE_DEPTH:
        raise FileToolError(f"Archive nesting exceeds {MAX_ARCHIVE_DEPTH} levels")
    entries = archive_listing(path)
    target = out / "extracted"
    target.mkdir(exist_ok=True)
    if zipfile.is_zipfile(path):
        with zipfile.ZipFile(path) as archive:
            archive.extractall(target)
    elif tarfile.is_tarfile(path):
        with tarfile.open(path) as archive:
            archive.extractall(target, filter="data")
    else:
        run(["7z", "x", "-y", f"-o{target}", str(path)])
    for child in target.rglob("*"):
        if child.is_symlink() or (child.exists() and not child.resolve().is_relative_to(target.resolve())):
            raise FileToolError("Archive extraction escaped its destination")
    listing = "\n".join(f"- {item['name']} ({item['size']} bytes)" for item in entries)
    return listing, listing, {"entries": entries, "expanded_bytes": sum(item["size"] for item in entries)}


def extract_legacy_office(path: Path, out: Path) -> tuple[str, str, dict[str, Any]]:
    profile = out / ".libreoffice-profile"
    converted = out / "converted"
    converted.mkdir(exist_ok=True)
    run(["libreoffice", "--headless", "--safe-mode", f"-env:UserInstallation=file://{profile}", "--convert-to", "pdf", "--outdir", str(converted), str(path)], cwd=out)
    pdfs = list(converted.glob("*.pdf"))
    if not pdfs:
        raise FileToolError("LibreOffice did not produce a PDF")
    text, markdown, metadata = extract_pdf(pdfs[0], out)
    metadata["converted_from"] = path.suffix.lower()
    return text, markdown, metadata


def extract(path: Path, out: Path) -> dict[str, Any]:
    inspection = inspect_file(path)
    if inspection.encrypted:
        raise FileToolError("Password-protected files are not supported")
    kind = inspection.category
    suffix = path.suffix.lower()
    if kind == "pdf":
        text, markdown, metadata = extract_pdf(path, out)
    elif suffix == ".docx":
        text, markdown, metadata = extract_docx(path)
    elif suffix == ".pptx":
        text, markdown, metadata = extract_pptx(path)
    elif suffix == ".xlsx":
        text, markdown, metadata = extract_xlsx(path)
    elif kind in {"document", "presentation", "spreadsheet"}:
        text, markdown, metadata = extract_legacy_office(path, out)
    elif kind == "ebook":
        text, markdown, metadata = extract_epub(path)
    elif kind == "email":
        text, markdown, metadata = extract_email(path)
    elif kind == "notebook":
        text, markdown, metadata = extract_notebook(path)
    elif kind == "image":
        text, markdown, metadata = extract_image(path, out)
    elif kind == "media":
        text, markdown, metadata = extract_media(path, out)
    elif kind == "archive":
        text, markdown, metadata = extract_archive(path, out)
    elif kind == "text":
        text = decode_text(path)
        markdown, metadata = text, {}
    else:
        raise FileToolError(f"Unsupported file format: {inspection.mime_type} ({suffix or 'no extension'})")
    encoded = text.encode("utf-8")
    truncated = len(encoded) > MAX_OUTPUT_BYTES
    if truncated:
        text = encoded[:MAX_OUTPUT_BYTES].decode("utf-8", errors="ignore")
        markdown = text
    (out / "content.txt").write_text(text, encoding="utf-8")
    (out / "content.md").write_text(markdown, encoding="utf-8")
    manifest = {**asdict(inspection), "metadata": metadata, "truncated": truncated, "outputs": {"text": "content.txt", "markdown": "content.md"}}
    (out / "manifest.json").write_text(json.dumps(manifest, ensure_ascii=False, indent=2), encoding="utf-8")
    return manifest


def self_test() -> dict[str, Any]:
    required = ["file", "pdftotext", "pdfinfo", "qpdf", "pandoc", "libreoffice", "ffmpeg", "ffprobe", "tesseract", "7z", "exiftool"]
    missing = [name for name in required if shutil.which(name) is None]
    modules = ["pymupdf", "PIL", "bs4", "docx", "pptx", "openpyxl", "ebooklib", "extract_msg", "nbformat", "magic", "yaml"]
    unavailable = []
    for name in modules:
        try:
            __import__(name)
        except Exception:
            unavailable.append(name)
    result = {"ok": not missing and not unavailable, "missing_binaries": missing, "missing_modules": unavailable}
    if not result["ok"]:
        raise FileToolError(json.dumps(result))
    return result


def main() -> int:
    parser = argparse.ArgumentParser(prog="notepatch-file")
    subparsers = parser.add_subparsers(dest="command", required=True)
    for name in ("inspect", "list"):
        command = subparsers.add_parser(name)
        command.add_argument("file")
    extract_parser = subparsers.add_parser("extract")
    extract_parser.add_argument("file")
    extract_parser.add_argument("--output-dir", required=True)
    subparsers.add_parser("self-test")
    args = parser.parse_args()
    try:
        if args.command == "self-test":
            result = self_test()
        else:
            path = source_path(args.file)
            if args.command == "inspect":
                result = asdict(inspect_file(path))
            elif args.command == "list":
                result = {"path": str(path), "entries": archive_listing(path)}
            else:
                result = extract(path, output_path(args.output_dir))
        print(json.dumps(result, ensure_ascii=False, indent=2))
        return 0
    except (FileToolError, OSError, ValueError, zipfile.BadZipFile, tarfile.TarError) as exc:
        print(json.dumps({"code": "file_parse_failed", "message": str(exc)}, ensure_ascii=False), file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
